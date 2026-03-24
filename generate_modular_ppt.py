from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation
prs = Presentation()

# Colors
ORANGE = RGBColor(249, 115, 22)
BLUE = RGBColor(56, 189, 248)
SLATE = RGBColor(51, 65, 85)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(16, 185, 129)

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    return slide

def draw_wide_flow_box(slide, text, top, fill_color=BLUE):
    left = Inches(1.0)
    width = Inches(8.0)
    height = Inches(0.8)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = SLATE
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def draw_centered_arrow(slide, top, height=0.4):
    left = Inches(4.8) # Center of a 10 inch slide is 5, but arrow is 0.4 wide
    width = Inches(0.4)
    height = Inches(height)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

def add_flow_slide(title_text, box_texts, final_green=False):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    title.text_frame.paragraphs[0].font.size = Pt(32)
    
    start_y = 1.6
    n = len(box_texts)
    
    if n == 0:
        return slide
    
    box_height = 0.8
    if n > 0:
        available_gap_space = 5.4 - (n * box_height)
        if n > 1:
            gap = max(0.4, available_gap_space / (n - 1))  # e.g., 5 boxes, 4 gaps
        else:
            gap = 0
            
    for i, text in enumerate(box_texts):
        # determine color
        color = GREEN if (final_green and i == n - 1) else BLUE
        
        y_pos = Inches(start_y + i * (box_height + gap))
        draw_wide_flow_box(slide, text, y_pos, fill_color=color)
        
        if i < n - 1:
            arrow_y = start_y + i * (box_height + gap) + box_height + 0.05
            arrow_height = max(0.2, gap - 0.1)
            draw_centered_arrow(slide, Inches(arrow_y), arrow_height)
            
    return slide

# --- SLIDES GENERATION ---

# Title
add_title_slide(
    "Tech Trolley System Process Flow",
    "Modular Diagrams for Operations & Exceptions"
)

# Fig 1
add_flow_slide(
    "Fig 1: Main Outward Flow (Godown to Venue)",
    [
        "1. Admin: Create Conference & Assign In-Charge Technician",
        "2. Tech: Login, View Assignment, Submit Requirements Digitally",
        "3. Godown: Receives List & Packs Materials into Cases",
        "4. System: Auto-Generates 3 Delivery Challans & E-Way Bills",
        "5. Logistics: Load Trucks, Transport, Setup at Venue"
    ]
)

# Fig 2
add_flow_slide(
    "Fig 2: Main Return Flow (Venue to Accounts)",
    [
        "1. Venue: Event Conclusion, Tech Counts & Verifies via App",
        "2. Office: Generate Return E-Way Bill (if > 25km)",
        "3. Logistics: Return Transport to Godown",
        "4. Godown: Final Strict Scan-In Count",
        "5. System & Accounts: Close Challan & Generate Final Invoice"
    ],
    final_green=True
)

# Fig 3.1
add_flow_slide(
    "Fig 3.1: Asset Transfer (Site-to-Site)",
    [
        "1. Tech A (Conference 1) initiates 'Asset Handshake'",
        "2. Tech B (Conference 2) Scans Assets",
        "3. System: Liability Shifts to Tech B",
        "4. Bypass Godown Entirely"
    ],
    final_green=True
)

# Fig 3.2
add_flow_slide(
    "Fig 3.2: Service & Defect Tracking",
    [
        "1. Tech Identifies Broken Item",
        "2. Tech Flags as 'Defective' in Mobile App",
        "3. System: Locks Asset in 'Maintenance Queue'",
        "4. Removed from Available Godown Pool"
    ]
)

# Fig 3.3
add_flow_slide(
    "Fig 3.3: End of Life / Scrapping",
    [
        "1. Identify Destroyed Asset",
        "2. Flag as 'Expired/Retired'",
        "3. Hidden from Godown Views",
        "4. Retained in Dashboard for ROI/Depreciation Logs"
    ]
)

# Fig 4.1
add_flow_slide(
    "Fig 4.1: Consumables Workflow",
    [
        "1. Scan Out: Worker Bypasses Scanner, Manually Types Qty",
        "2. Packing Up: Venue Supervisor Types Remaining Qty",
        "3. Inwarding: Godown Confirms Leftover Amount",
        "4. System Calculates Exact Consumption for Billing"
    ],
    final_green=True
)

# Fig 4.2
add_flow_slide(
    "Fig 4.2: Demo & Internal Services",
    [
        "1. Start: Define as 'Internal' Conference",
        "2. System: Generates Zero-Value Delivery Challan",
        "3. Gate Pass: Works Normally (Security satisfied)",
        "4. Finance: Bypassed strictly for Accounts"
    ],
    final_green=True
)

# Fig 4.3
add_flow_slide(
    "Fig 4.3: External Rental (3rd Party)",
    [
        "1. Admin Manually Types Items in Rental Asset List",
        "2. Rented Items get printed on Delivery Challan",
        "3. Items Bypass the Godown QR Crosscheck"
    ],
    final_green=True
)

# Fig 4.4
add_flow_slide(
    "Fig 4.4: Personal Assets",
    [
        "1. Technician Requests Personal Asset from Inventory",
        "2. Admin Marks Asset as 'Taken' along with Technician's Name",
        "3. Asset Status Tracked in Inventory Audit List"
    ]
)

out_file = "Modular_Process_Flowcharts.pptx"
prs.save(out_file)
print(f"Presentation successfully built at: {out_file}")
