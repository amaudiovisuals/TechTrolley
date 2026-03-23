from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation
prs = Presentation()

# Colors
ORANGE = RGBColor(249, 115, 22)
BLUE = RGBColor(56, 189, 248)
SLATE = RGBColor(51, 65, 85)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = ORANGE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    return slide

def add_bullet_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.color.rgb = ORANGE
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    tf = body_shape.text_frame
    
    if len(bullet_points) > 0:
        tf.paragraphs[0].text = bullet_points[0]
        tf.paragraphs[0].font.size = Pt(16)
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            if point.startswith("  -") or point.startswith("    "):
                p.level = 1
                p.font.size = Pt(14)
            else:
                p.font.size = Pt(16)
            
    return slide

def draw_wide_flow_box(slide, text, top, fill_color=BLUE):
    # Centered, very wide box to ensure text never overflows
    left = Inches(1.0)
    width = Inches(8.0)
    height = Inches(0.8)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = SLATE
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def draw_centered_arrow(slide, top):
    # Downward arrow perfectly between boxes
    left = Inches(4.8) # Center of a 10 inch slide is 5, but arrow is 0.4 wide
    width = Inches(0.4)
    height = Inches(0.4)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

# --- SLIDES GENERATION ---

# Slide 1: Intro
add_title_slide(
    "Tech Trolley Platform",
    "Logistics, Asset Tracking, & Process Revolution"
)

# Slide 2: Objectives
add_bullet_slide(
    "Platform Objectives & Core Features",
    [
        "Digitizing legacy hardware tracking into a unified SaaS platform.",
        "  - Total Accountability: Track, maintain, and control inventory.",
        "  - Automation: QR/Barcode generation & instant scanning.",
        "  - Logistics Engine: Auto-generating Delivery Challans.",
        "  - Double Verification: Scanning rules for Godown <-> Venue.",
        "  - Deep Analytics: Track asset usage history and depreciation."
    ]
)

# Slide 3: Flow 1 (Outward)
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
slide3.shapes.title.text = "Process Flow 1: Conference Events (Outward)"
slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = ORANGE

start_y = Inches(1.5)
gap = Inches(1.3)

draw_wide_flow_box(slide3, "1. Technician places requirements digitally via the App", start_y)
draw_centered_arrow(slide3, start_y + Inches(0.85))

draw_wide_flow_box(slide3, "2. Godown In-Charge receives list & packs materials into Cases", start_y + gap)
draw_centered_arrow(slide3, start_y + gap + Inches(0.85))

draw_wide_flow_box(slide3, "3. System Auto-Generates 3 Delivery Challans & E-Way Bills", start_y + (gap*2))
draw_centered_arrow(slide3, start_y + (gap*2) + Inches(0.85))

draw_wide_flow_box(slide3, "4. Materials Loaded onto Trucks, Transported, & Setup at Venue", start_y + (gap*3))

# Slide 4: Flow 1 (Return)
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
slide4.shapes.title.text = "Process Flow 1: Conference Events (Return)"
slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = ORANGE

draw_wide_flow_box(slide4, "5. Event Over: Technician counts & verifies materials via Mobile App", start_y)
draw_centered_arrow(slide4, start_y + Inches(0.85))

draw_wide_flow_box(slide4, "6. Office generates Return E-Way Bill (if venue > 25km)", start_y + gap)
draw_centered_arrow(slide4, start_y + gap + Inches(0.85))

draw_wide_flow_box(slide4, "7. Return to Godown: Godown In-Charge does final strict scan-in count", start_y + (gap*2))
draw_centered_arrow(slide4, start_y + (gap*2) + Inches(0.85))

draw_wide_flow_box(slide4, "8. Delivery Challan Officially Closed & Sent to Accounts for Invoice", start_y + (gap*3), fill_color=RGBColor(16, 185, 129))

# Slide 5: Flow 2
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
slide5.shapes.title.text = "Process Flow 2: Direct Rental"
slide5.shapes.title.text_frame.paragraphs[0].font.color.rgb = ORANGE

draw_wide_flow_box(slide5, "1. Technician Places Rental Requirement", start_y)
draw_centered_arrow(slide5, start_y + Inches(0.85))

draw_wide_flow_box(slide5, "2. Godown packs items & Generates Direct Delivery Challans", start_y + gap)
draw_centered_arrow(slide5, start_y + gap + Inches(0.85))

draw_wide_flow_box(slide5, "3. Assets handed directly to Customer at Godown", start_y + (gap*2))
draw_centered_arrow(slide5, start_y + (gap*2) + Inches(0.85))

draw_wide_flow_box(slide5, "4. Return Condition Inspected => Invoice Generated", start_y + (gap*3), fill_color=RGBColor(16, 185, 129))

# Slide 6: Exceptions
add_bullet_slide(
    "Exception Solution: Consumables Workflow",
    [
        "Problem: We cannot use QR codes to track how much battery or tape comes back.",
        "System Implementation:",
        "  - The Inventory Engine will get a 'Consumable' toggle switch.",
        "  - When Scanning Out: Workers bypass the scanner and simply type the Quantity taken.",
        "  - When Packing Up: The Venue Supervisor types the exact remaining Quantity left.",
        "  - When Inwarding: Godown confirms the leftover amount. The system instantly calculates the exact consumption and routes it to billing."
    ]
)

# Slide 7: Site to Site
add_bullet_slide(
    "Exception Solutions: Logistics & Service",
    [
        "Site-to-Site Transfer",
        "  - Problem: Materials moving directly between venues.",
        "  - Implementation: The App allows an 'Asset Handshake'. Tech A scans gear mapping it to 'Conference B'. The system shifts liability to Tech B, bypassing Godown entirely.",
        "",
        "Service & Defect Tracking",
        "  - Problem: Broken materials on-site.",
        "  - Implementation: Technicians hit 'Flag Defective' during return scanning. The asset is removed from the Godown Available Pool and locked inside a 'Maintenance Queue'."
    ]
)

# Slide 8: End of Life
add_bullet_slide(
    "Exception Solutions: Loans & EOL",
    [
        "Demo & Third-Party Services",
        "  - Implementation: We create 'Internal' conferences in the database. This allows Godown to generate standard Delivery Challans (zero value) so Gate passes still work flawlessly without alerting financial accounts.",
        "",
        "End of Life / Scrapping",
        "  - Problem: We need to remove destroyed items without losing historical invoices.",
        "  - Implementation: Assets are flagged 'Retired'. They disappear from Godown views but remain strictly in the Web Dashboard's historical logs for Depreciation and ROI math."
    ]
)

add_title_slide("Ready for Deployment", "Phase 2 Commences Now.")

prs.save("Tech_Trolley_Presentation.pptx")
print("Presentation successfully built!")
